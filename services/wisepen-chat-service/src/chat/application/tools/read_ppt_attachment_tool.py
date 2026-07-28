from typing import Dict, Any

import httpx

from chat.core.config.app_settings import settings
from chat.application.tools.core import (
    ToolDefinition,
    ToolExecutionError,
    ToolLLMSpec,
    ToolParametersSchema,
    ToolPolicy,
    ToolRiskLevel,
)
from chat.domain.repositories import SessionRepository
from chat.service_client.file_storage_service_client import FileStorageClient


class ReadPptAttachmentTool:
    """Extract text content from a PowerPoint (.pptx) attachment."""

    def __init__(self, file_storage_client: FileStorageClient, session_repo: SessionRepository) -> None:
        self._file_storage = file_storage_client
        self._session_repo = session_repo
        parameters_schema: Dict[str, Any] = {
            "type": "object",
            "properties": {
                "attachment_id": {
                    "type": "string",
                    "description": "The attachment ID to read.",
                },
            },
            "required": ["attachment_id"],
        }
        self._definition = ToolDefinition(
            llm_spec=ToolLLMSpec(
                name="read_ppt_attachment",
                description=(
                    "Extract and read the text content of a PowerPoint (.pptx) attachment. "
                    "Returns slide-by-slide text. Use this for .pptx files attached to the conversation."
                ),
                parameters_schema=ToolParametersSchema(parameters_schema),
            ),
            policy=ToolPolicy(
                expose_by_default=True,
                persist_output=False,
                risk_level=ToolRiskLevel.LOW,
                required_context_keys=("user_id", "session_id"),
                timeout_seconds=30.0,
                max_output_chars=settings.TOOL_RESULT_MAX_CHARS,
            ),
        )

    @property
    def definition(self) -> ToolDefinition:
        return self._definition

    async def execute(
        self,
        context: dict[str, Any],
        config: dict[str, Any] | None = None,
        **kwargs: Any,
    ) -> str:
        user_id = context.get("user_id")
        session_id = context.get("session_id")
        attachment_id = kwargs.get("attachment_id", "").strip()
        if not attachment_id:
            raise ToolExecutionError(reason="missing_attachment_id", detail_reason="Missing required argument: attachment_id.")

        temp_refs, _res_refs = await self._session_repo.get_session_attachments(session_id, user_id)
        ref = next((r for r in (temp_refs or []) if r.attachment_id == attachment_id), None)
        if ref is None:
            raise ToolExecutionError(reason="attachment_not_found", detail_reason=f"Attachment {attachment_id} not found in session.")

        download_url = await self._file_storage.get_download_url(ref.object_key)
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                resp = await client.get(download_url)
                resp.raise_for_status()
                pptx_bytes = resp.content
        except httpx.HTTPError as e:
            raise ToolExecutionError(reason="download_failed", detail_reason=f"Failed to download attachment: {e}") from e

        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(pptx_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                shapes_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shapes_text.append(shape.text_frame.text.strip())
                slide_content = "\n".join(shapes_text)
                if slide_content:
                    slides_text.append(f"--- Slide {i} ---\n{slide_content}")
            text = "\n\n".join(slides_text)
        except ImportError:
            raise ToolExecutionError(reason="pptx_parser_unavailable", detail_reason="PowerPoint parser (python-pptx) is not installed.")
        except Exception as e:
            raise ToolExecutionError(reason="pptx_parse_failed", detail_reason=f"Failed to parse PowerPoint: {e}") from e

        if len(text) > settings.TOOL_RESULT_MAX_CHARS:
            text = text[:settings.TOOL_RESULT_MAX_CHARS] + "\n...[truncated]"
        return text
