from io import BytesIO
from typing import Callable, Dict, List, Optional
from zipfile import BadZipFile, ZipFile

import pdfplumber
from docx import Document as DocxDocument
from markdownify import markdownify
from openpyxl import load_workbook
from pptx import Presentation
from readability import Document

from common.logger import log_ok, log_fail

__all__ = [
    "ContentProcessor",
    "DocumentParser",
]

OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

MAX_ZIP_UNCOMPRESSED_SIZE = 200 * 1024 * 1024
ANTI_CRAWL_SCAN_CHARS = 20000

ANTI_CRAWL_KEYWORDS = (
    "just a moment",
    "please enable javascript",
    "please enable js",
    "captcha",
    "access denied",
    "cloudflare",
    "are you a robot",
    "are you human",
    "verify you are human",
)


def looks_like_anti_crawl(text: str) -> bool:
    lower = text[:ANTI_CRAWL_SCAN_CHARS].lower()
    return any(kw in lower for kw in ANTI_CRAWL_KEYWORDS)


def normalize_text(text: str) -> str:
    lines = [
        line.rstrip()
        for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    ]
    result = "\n".join(lines).strip()

    while "\n\n\n" in result:
        result = result.replace("\n\n\n", "\n\n")

    return result


class DocumentParser:
    """文档解析器：负责 bytes 文档识别、解析和基础质量过滤。"""

    def __init__(
        self,
        min_content_length: int = 50,
        max_document_size: int = 50 * 1024 * 1024,
    ):
        self._min_content_length = min_content_length
        self._max_document_size = max_document_size
        self._parsers: Dict[str, Callable[[bytes], Optional[str]]] = {
            "pdf": self._parse_pdf,
            "docx": self._parse_docx,
            "xlsx": self._parse_xlsx,
            "pptx": self._parse_pptx,
        }

    def parse(self, data: bytes) -> Optional[str]:
        if len(data) > self._max_document_size:
            log_fail("文档解析", f"文件过大({len(data)}字节)，上限{self._max_document_size}字节")
            return None

        doc_type = self._detect_doc_type(data)
        if doc_type is None:
            return None

        parser = self._parsers.get(doc_type)
        if parser is None:
            log_fail("文档解析", f"暂不支持的文档类型: {doc_type}")
            return None

        text = parser(data)
        if text is None:
            log_fail("文档清洗", "文本提取返回空")
            return None

        cleaned = normalize_text(text)

        if looks_like_anti_crawl(cleaned):
            log_fail("文档清洗", "提取文本疑似反爬/错误页面，触发降级")
            return None

        if len(cleaned) < self._min_content_length:
            log_fail(
                "文档清洗",
                f"提取文本过短({len(cleaned)}字符)，阈值{self._min_content_length}，触发降级",
            )
            return None

        log_ok("文档清洗", length=len(cleaned), doc_type=doc_type)
        return cleaned

    def _detect_doc_type(self, data: bytes) -> Optional[str]:
        if data[:5] == b"%PDF-":
            return "pdf"

        if data[:8] == OLE_MAGIC:
            log_fail("文档解析", "OLE 复合文档(旧版 .doc/.xls/.ppt)，暂不支持")
            return None

        try:
            with ZipFile(BytesIO(data)) as zf:
                infos = zf.infolist()
                uncompressed_size = sum(info.file_size for info in infos)

                if uncompressed_size > MAX_ZIP_UNCOMPRESSED_SIZE:
                    log_fail(
                        "文档解析",
                        f"ZIP 解压后体积过大({uncompressed_size}字节)，上限{MAX_ZIP_UNCOMPRESSED_SIZE}字节",
                    )
                    return None

                names = {info.filename for info in infos}

                matches = [
                    ("word/document.xml", "docx"),
                    ("xl/workbook.xml", "xlsx"),
                    ("ppt/presentation.xml", "pptx"),
                ]

                hits = [doc_type for sig, doc_type in matches if sig in names]

                if len(hits) != 1:
                    if hits:
                        log_fail("文档解析", f"ZIP 内含多种 Office 特征文件({', '.join(hits)})，无法确定类型")
                    return None

                return hits[0]

        except BadZipFile:
            pass

        log_fail("文档解析", "无法识别文档类型")
        return None

    def _parse_pdf(self, data: bytes) -> Optional[str]:
        try:
            text_parts: List[str] = []
            page_count = 0

            with pdfplumber.open(BytesIO(data)) as pdf:
                page_count = len(pdf.pages)

                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            result = "\n".join(text_parts).strip()

            if result:
                log_ok("PDF 文本提取", pages=page_count, length=len(result))
                return result

            return None

        except Exception as e:
            log_fail("PDF 文本提取", e)
            return None

    def _parse_docx(self, data: bytes) -> Optional[str]:
        try:
            doc = DocxDocument(BytesIO(data))
            parts: List[str] = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))

            result = "\n".join(parts).strip()

            if result:
                log_ok("DOCX 文本提取", length=len(result))
                return result

            return None

        except Exception as e:
            log_fail("DOCX 文本提取", e)
            return None

    def _parse_xlsx(self, data: bytes) -> Optional[str]:
        wb = None

        try:
            wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
            parts: List[str] = []

            for sheet in wb:
                sheet_rows: List[str] = []

                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    line = "\t".join(cells).strip()

                    if line:
                        sheet_rows.append(line)

                if sheet_rows:
                    parts.append(f"## Sheet: {sheet.title}")
                    parts.append("")
                    parts.append("```tsv")
                    parts.extend(sheet_rows)
                    parts.append("```")
                    parts.append("")

            result = "\n".join(parts).strip()

            if result:
                log_ok("XLSX 文本提取", length=len(result))
                return result

            return None

        except Exception as e:
            log_fail("XLSX 文本提取", e)
            return None

        finally:
            if wb is not None:
                wb.close()

    def _parse_pptx(self, data: bytes) -> Optional[str]:
        try:
            prs = Presentation(BytesIO(data))
            parts: List[str] = []

            for index, slide in enumerate(prs.slides, 1):
                slide_parts: List[str] = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_parts.append(text)

                if slide_parts:
                    parts.append(f"## Slide {index}")
                    parts.append("")
                    parts.extend(slide_parts)
                    parts.append("")

            result = "\n".join(parts).strip()

            if result:
                log_ok("PPTX 文本提取", length=len(result))
                return result

            return None

        except Exception as e:
            log_fail("PPTX 文本提取", e)
            return None


class ContentProcessor:
    """内容处理器：负责 str/bytes 分流，以及 HTML/纯文本清洗。"""

    def __init__(
        self,
        min_content_length: int = 400,
        document_min_content_length: int = 50,
        max_document_size: int = 50 * 1024 * 1024,
    ):
        self._min_content_length = min_content_length
        self._document_parser = DocumentParser(
            min_content_length=document_min_content_length,
            max_document_size=max_document_size,
        )

    def process(self, content: str | bytes) -> Optional[str]:
        if isinstance(content, bytes):
            return self._document_parser.parse(content)

        return self._process_text(content)

    def _process_text(self, content: str) -> Optional[str]:
        stripped = content.strip()
        if not stripped:
            return None

        if looks_like_anti_crawl(stripped):
            log_fail("内容检测", "疑似反爬/错误页面，触发降级")
            return None

        lower_head = stripped[:1024].lower()
        if "<html" in lower_head or "<!doctype html" in lower_head or "<body" in lower_head:
            return self._process_html(stripped)

        return self._process_plain_text(stripped)

    def _process_html(self, html: str) -> Optional[str]:
        try:
            clean_content = self._extract_main_content(html)
            markdown = self._convert_to_markdown(clean_content)
            result = normalize_text(markdown)

            if looks_like_anti_crawl(result):
                log_fail("HTML 清洗", "清洗后疑似反爬/错误页面，触发降级")
                return None

            if len(result) < self._min_content_length:
                log_fail(
                    "HTML 清洗",
                    f"清洗后文本过短({len(result)}字符)，阈值{self._min_content_length}，触发降级",
                )
                return None

            return result

        except Exception as e:
            log_fail("HTML 清洗", e, fallback="返回未清洗原文，可能含 HTML 标签")
            fallback = normalize_text(html)

            if len(fallback) < self._min_content_length:
                return None

            return fallback

    def _process_plain_text(self, text: str) -> Optional[str]:
        normalized = normalize_text(text)

        if len(normalized) < self._min_content_length:
            log_fail(
                "纯文本检测",
                f"文本过短({len(normalized)}字符)，阈值{self._min_content_length}，触发降级",
            )
            return None

        if looks_like_anti_crawl(normalized):
            log_fail("纯文本检测", "疑似反爬/错误页面，触发降级")
            return None

        return normalized

    def _extract_main_content(self, html: str) -> str:
        try:
            return Document(html).summary()
        except Exception as e:
            log_fail("readability 提取主体内容", e)
            return html

    def _convert_to_markdown(self, html: str) -> str:
        try:
            return markdownify(
                html,
                heading_style="ATX",
                strip=["script", "style", "img", "nav", "footer"],
                autolinks=True,
            ).strip()
        except Exception as e:
            log_fail("HTML 转 Markdown", e)
            return html
